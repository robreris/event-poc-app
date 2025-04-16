from pptx import Presentation
from pathlib import Path

NOTES_DIR = Path("notes")
NOTES_DIR.mkdir(exist_ok=True)

def extract_notes(pptx_path: Path) -> list:
    prs = Presentation(pptx_path)
    output = []

    for i, slide in enumerate(prs.slides, start=1):
        notes_slide = slide.notes_slide if slide.has_notes_slide else None
        notes_text = notes_slide.notes_text_frame.text if notes_slide else ""

        slide_info = {
            "slide": i,
            "notes": notes_text.strip()
        }

        # Write to individual text file
        note_path = NOTES_DIR / f"slide_{i:02d}.txt"
        with open(note_path, "w", encoding="utf-8") as f:
            f.write(notes_text)

        output.append(slide_info)

    return output
