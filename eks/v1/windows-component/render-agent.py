import os
import time
import json
from pathlib import Path
import boto3
import win32com.client
import pika
from pptx import Presentation

# --- RabbitMQ Config from Environment ---
RABBITMQ_HOST = os.getenv("RABBITMQ_HOST", "a1aa033d6ebf746c69dcec222772a47e-1819495256.us-east-1.elb.amazonaws.com")
RABBITMQ_USER = os.getenv("RABBITMQ_USER", "default_user_4QGwTp4o-uNzyq6G0ki")
RABBITMQ_PASS = os.getenv("RABBITMQ_PASS", "jkyzYCrd--jlZKXWhbn__Rk9aELWpShr")
RABBITMQ_PORT = int(os.getenv("RABBITMQ_PORT", "5672"))
QUEUE_NAME = os.getenv("RABBITMQ_QUEUE", "ppt_tasks")
RESPONSE_QUEUE = os.getenv("RABBITMQ_RESPONSE_QUEUE", "windows_response")

# === Config from environment ===
S3_BUCKET = os.getenv("S3_BUCKET", "event-driven-poc-ftnt")
S3_REGION = os.getenv("S3_REGION", "us-east-1")
TMP_DIR = Path(os.getenv("TMP_DIR", "C:/Temp/pptx-jobs"))
TMP_DIR.mkdir(parents=True, exist_ok=True)

# === S3 Utilities ===

def download_from_s3(s3_key, local_path):
    s3 = boto3.client(
        "s3",
        region_name=S3_REGION,
    )
    s3.download_file(S3_BUCKET, s3_key, str(local_path))

def upload_to_s3(local_path, s3_key):
    s3 = boto3.client("s3", region_name=S3_REGION)
    s3.upload_file(str(local_path), S3_BUCKET, s3_key)
    print(f"[INFO] Uploaded {local_path} to s3://{S3_BUCKET}/{s3_key}")
    return f"s3://{S3_BUCKET}/{s3_key}"

# === PowerPoint Extraction ===

### Note: need to work in functionality to extract notes and upload to their own folder in S3

def process_pptx_notes(pptx_path, output_dir):
    prs = Presentation(pptx_path)
    output = []

    for i, slide in enumerate(prs.slides, start=1):
        notes_slide = slide.notes_slide if slide.has_notes_slide else None
        notes_text = notes_slide.notes_text_frame.text if notes_slide else ""
        note_path = output_dir / f"slide_{i:02d}.txt"
        
        slide_info = {
            "slide": i,
            "notes": notes_text.strip()
        }

        # Write to individual text file      
        with open(note_path, "w", encoding="utf-8") as f:
            f.write(notes_text)

        output.append(slide_info)

    return output

# Wait for the file to exist and grow stable in size
def wait_for_video(path, min_size=10000, stable_seconds=3, timeout=120):
    """Wait for video to be written, stable for several seconds."""
    start = time.time()
    last_size = -1
    stable_time = None
    while True:
        if os.path.exists(path):
            size = os.path.getsize(path)
            if size >= min_size:
                if size == last_size:
                    if stable_time is None:
                        stable_time = time.time()
                    elif time.time() - stable_time >= stable_seconds:
                        return True
                else:
                    stable_time = None
                last_size = size
        if time.time() - start > timeout:
            print(f"[WARN] Timeout waiting for {path} to stabilize.")
            return False
        time.sleep(1)

def process_pptx_slides(pptx_path, output_dir):
    ppt_app = win32com.client.Dispatch("PowerPoint.Application")
    ppt_app.Visible = 1
    pres = ppt_app.Presentations.Open(str(pptx_path), WithWindow=False)

    slides_info = []
    videos_info = []

    output_dir = Path(output_dir)
    output_dir.mkdir(exist_ok=True, parents=True)

    for idx, slide in enumerate(pres.Slides):
        slide_number = idx + 1
        slide_id = slide.SlideID
        try:
            slide_title = slide.Shapes.Title.TextFrame.TextRange.Text
        except Exception:
            slide_title = f"Slide {slide_number}"

        # Detect animation: presence of effects in MainSequence
        has_animation = slide.TimeLine.MainSequence.Count > 0

        if has_animation:
            # Export animation as mp4 using CreateVideo workaround
            video_filename = f"slide_{slide_number:02d}_animation.mp4"
            video_path = str(output_dir / video_filename)
            
            try:
                # Create a new temp presentation with just this slide
                temp_pres = ppt_app.Presentations.Add()
                while temp_pres.Slides.Count > 0:
                    temp_pres.Slides[1].Delete()
                slide.Copy()
                temp_pres.Slides.Paste()
                temp_pres.SaveAs(str(output_dir / f"slide_{slide_number:02d}_only.pptx"))
                for _ in range(5):
                    if Path(video_path).exists():
                        break
                    time.sleep(1)
                time.sleep(0.5)  # Give a little extra time for good measure
                temp_pres.CreateVideo(video_path, UseTimingsAndNarrations=True)
                print(f"CreateVideoStatus for slide_{slide_number:02d}_only.pptx: {temp_pres.CreateVideoStatus}")

                print(f"[INFO] Exporting animation for slide {slide_number} to {video_path}...")
                # Wait for export to finish (poll CreateVideoStatus)
                while temp_pres.CreateVideoStatus == 1:  # 1 == InProgress, 2 == Done
                    time.sleep(1)
                    
                for _ in range(20):
                    if os.path.exists(video_path) and os.path.getsize(video_path) > 10000:
                        break
                    time.sleep(1)
                # Now wait for the file to stabilize:
                wait_for_video(video_path, min_size=500, stable_seconds=5, timeout=180)
                temp_pres.Close()
            except Exception as error:
                print(f"Error exporting animation: {error}")

            videos_info.append({
                "slide_id": slide_id,
                "slide_number": slide_number,
                "s3_key": None,  # to be filled after upload
                "title": slide_title,
                "type": "animation",
                "local_path": video_path
            })
        else:
            # Export static slide as PNG
            img_filename = f"slide_{slide_number:02d}.png"
            img_path = str(output_dir / img_filename)
            slide.Export(img_path, "PNG")
            slides_info.append({
                "slide_id": slide_id,
                "slide_number": slide_number,
                "s3_key": None,  # to be filled after upload
                "title": slide_title,
                "type": "static",
                "local_path": img_path
            })

    pres.Close()
    ppt_app.Quit()
    return slides_info, videos_info

# === Job Processing Logic ===

def process_job(job):
    job_id = job["job_id"]
    pptx_s3_key = job["s3_key"]
    pptx_filename = job["filename"]
    local_pptx = TMP_DIR / f"{job_id}_{pptx_filename}"
    slides_dir = TMP_DIR / f"{job_id}_slides"
    slides_dir.mkdir(exist_ok=True)
    notes_dir = TMP_DIR / f"{job_id}_notes"
    notes_dir.mkdir(exist_ok=True)

    print(f"[INFO] Downloading {pptx_s3_key} from S3 to {local_pptx}")
    download_from_s3(pptx_s3_key, local_pptx)
    print("[INFO] Download complete. Processing slides/animations...")

    notes = process_pptx_notes(local_pptx, notes_dir)
    slides, videos = process_pptx_slides(local_pptx, slides_dir)

    # Upload notes
    notes_files = []
    for text_file in sorted(notes_dir.glob("*.txt")):
        local_path = notes_dir / text_file
        s3_key = f"processed/notes/{job_id}/{Path(local_path).name}"
        upload_to_s3(local_path, s3_key)
        notes_files.append(Path(local_path).name)
        del local_path

    # Upload each exported file to S3
    slides_files = []
    for obj in slides + videos:
        local_path = obj["local_path"]
        s3_key = f"processed/slides/{job_id}/{Path(local_path).name}"
        upload_to_s3(local_path, s3_key)
        slides_files.append(Path(local_path).name)
        del obj["local_path"]  # remove local path, not needed for result
    
    tts_engine = job["tts_engine"]
    piper_args = job["piper_args"]
    file_id = job["file_id"]
    
    # Compose manifest JSON
    result = {
        "job_id": job_id,
        "s3_bucket": S3_BUCKET,
        "notes_s3_prefix": f"processed/notes/{job_id}",
        "notes_files": notes_files,
        "slides_s3_prefix": f"processed/slides/{job_id}",
        "slides_files": slides_files,
        "voice": job["voice"],
        "tts_engine": tts_engine,
        "piper_args": piper_args,
        "file_id": file_id
        
    }
    return result

# === Local Test Harness ===

def test_main():
    # Simulated incoming job (edit as needed)
    test_message = {
        "job_id": "job123",
        "s3_key": "Slides for FGT FSI Script - Copy.pptx",
        "filename": "sample.pptx",
        "efs_uploads": [
            {"filename": "user_video_1.mp4", "efs_path": "/mnt/efs/user_abc/user_video_1.mp4"}
        ]
    }
    result = process_job(test_message)
    manifest_path = TMP_DIR / f"{test_message['job_id']}_manifest.json"
    with open(manifest_path, "w") as f:
        json.dump(result, f, indent=2)
    print(f"\n[INFO] Job complete. Manifest written to {manifest_path}")
    print(json.dumps(result, indent=2))

def publish_response(message):
    credentials = pika.PlainCredentials(RABBITMQ_USER, RABBITMQ_PASS)
    params = pika.ConnectionParameters(
        host=RABBITMQ_HOST,
        port=RABBITMQ_PORT,
        credentials=credentials
    )
    with pika.BlockingConnection(params) as connection:
        channel = connection.channel()
        channel.queue_declare(queue=RESPONSE_QUEUE, durable=True)
        channel.basic_publish(
            exchange='',
            routing_key=RESPONSE_QUEUE,
            body=message,
            properties=pika.BasicProperties(delivery_mode=2)
        )
        print(f"[INFO] Published message to {RESPONSE_QUEUE}")


# === RabbitMQ Functions
def rabbitmq_callback(ch, method, properties, body):
    try:
        print("[INFO] Received message from RabbitMQ.")
        job = json.loads(body)
        result = process_job(job)
        print("TYPE BEFORE JSON:", type(result))
        print("RESULT CONTENTS:", result)
        # Publish result as a JSON string to the response queue
        #ch.basic_publish(
        #    exchange='',
        #    routing_key=RESPONSE_QUEUE,
        #    body=json.dumps(result)
        #)
        publish_response(json.dumps(result))
        print(f"[INFO] Sent result for job {job['job_id']} to {RESPONSE_QUEUE}")
    except Exception as e:
        print(f"[ERROR] {e}")
        ch.basic_publish(
            exchange='',
            routing_key=RESPONSE_QUEUE,
            body=json.dumps({"status": "failed", "error": str(e)})
        )
    ch.basic_ack(delivery_tag=method.delivery_tag)

def main_rabbitmq():
    while True:
        try:
            credentials = pika.PlainCredentials(RABBITMQ_USER, RABBITMQ_PASS)
            params = pika.ConnectionParameters(
                host=RABBITMQ_HOST,
                port=RABBITMQ_PORT,
                credentials=credentials,
                heartbeat=600,
                blocked_connection_timeout=300,
                connection_attempts=3,
                retry_delay=5
            )
            connection = pika.BlockingConnection(params)
            channel = connection.channel()
            channel.queue_declare(queue=QUEUE_NAME, durable=True)
            channel.queue_declare(queue=RESPONSE_QUEUE, durable=True)
            channel.basic_qos(prefetch_count=1)
            channel.basic_consume(queue=QUEUE_NAME, on_message_callback=rabbitmq_callback)
            print(f"[INFO] Awaiting jobs on queue '{QUEUE_NAME}' ...")
            channel.start_consuming()
        except Exception as e:
            print(f"[WARN] Could not connect to RabbitMQ: {e}. Retrying in 10s...")
            time.sleep(10)

if __name__ == "__main__":
    import sys
    # Usage: script.py [rabbitmq]
    if len(sys.argv) > 1 and sys.argv[1] == "rabbitmq":
        main_rabbitmq()
    else:
        test_main()

